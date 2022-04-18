from pptx import Presentation
from pptx.util import Inches, Pt

ppt = Presentation()


blank_slide_layout = ppt.slide_layouts[6]

slide = ppt.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)

txBox = slide.shapes.add_textbox(left, top,
                                 width, height)

tf = txBox.text_frame
tf.text = "This is text inside a textbox"

p = tf.add_paragraph()

p.text = "This is a second paragraph that's bold and italic"

# font
p.font.bold = True
p.font.italic = True

p = tf.add_paragraph()
p.text = "This is a third paragraph that's big "
p.font.size = Pt(40)

# save file
ppt.save('test_2.pptx')

print("done")